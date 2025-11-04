from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from backend.core.llm_clients import llm_client
from datetime import datetime
import re, os


class PPTAgent:
    def __init__(self, faiss_store):
        self.faiss_store = faiss_store
        self.name = "PPT Creation Agent"

    def clean_text(self, text: str) -> str:
        """Remove markdown symbols and unwanted characters."""
        text = re.sub(r'^#{1,6}\s*', '', text, flags=re.MULTILINE)
        text = re.sub(r"[*#`_~]+", "", text)
        text = re.sub(r"\s+", " ", text)
        text = re.sub(r"^[-*•]\s*", "", text.strip())
        return text.strip()

    async def process(self, query: str, pdf_text: str) -> str:
        """Generate a clean 3-slide PowerPoint presentation."""
        
        system_prompt =  """Create exactly 3 slides. Format strictly as:

SLIDE 1
Title (max 7 words)
Line 1 (around 12 words)
Line 2 (around 12 words)
Line 3 (around 12 words)

SLIDE 2
Title (max 7 words)
Line 1 (around 12 words)
Line 2 (around 12 words)
Line 3 (around 12 words)

SLIDE 3
Title (max 7 words)
Line 1 (around 12 words)
Line 2 (around 12 words)
Line 3 (around 12 words)

NO markdown, NO symbols, NO dashes. Just plain text.
SLIDE 1 must introduce what the PDF is about with apt title
Each line must be concise (around 12 words).
First slide should introduce what the PDF is about.
Slide 2 and 3 should summarize other key topics or insights from the PDF with apt title.
Keep the output clean, short, and fast to generate."""


        prompt = f"{query}\n\nDocument:\n{pdf_text[:3000]}"

        # Get AI response
        response = await llm_client.generate(prompt, system_prompt)

        # Parse slides by splitting on "SLIDE"
        slide_parts = re.split(r'SLIDE\s+\d+', response, flags=re.IGNORECASE)
        slide_parts = [s.strip() for s in slide_parts if s.strip()]

        if len(slide_parts) < 3:
            return " Failed to generate 3 slides. Try again."

        prs = Presentation()

        # Create slides
        for idx in range(min(3, len(slide_parts))):
            block = slide_parts[idx]
            lines = [self.clean_text(ln) for ln in block.split('\n') if ln.strip()]
            
            if len(lines) < 2:
                continue
            
            # First line = title, rest = bullets
            title = lines[0][:55]
            bullets = [ln for ln in lines[1:5] if ln]
            
            if not bullets:
                continue
            
            # Use same layout for ALL slides (layout 1 = title + content)
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            
            # Title styling
            slide.shapes.title.text = title
            title_frame = slide.shapes.title.text_frame.paragraphs[0]
            title_frame.font.size = Pt(34)
            title_frame.font.bold = True
            title_frame.font.color.rgb = RGBColor(139, 0, 0)
            
            # Content - check if placeholder exists
            if len(slide.placeholders) > 1:
                content = slide.placeholders[1]
                tf = content.text_frame
                tf.clear()
                
                for bullet in bullets:
                    p = tf.add_paragraph()
                    p.text = bullet
                    p.font.size = Pt(26)
                    p.font.color.rgb = RGBColor(0, 80, 80)
                    p.level = 0
            else:
                # Fallback if no content placeholder
                left, top = Inches(1), Inches(2.5)
                width, height = Inches(8), Inches(4)
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                
                for bullet in bullets:
                    p = tf.add_paragraph()
                    p.text = f"• {bullet}"
                    p.font.size = Pt(18)
                    p.font.color.rgb = RGBColor(0, 80, 80)

        # Save
        os.makedirs("outputs", exist_ok=True)
        filename = f"generated_ppt_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        output_path = os.path.join("outputs", filename)
        prs.save(output_path)

        return output_path
